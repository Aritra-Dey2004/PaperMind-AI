from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import json
import re


def extract_json(text):
    text = text.strip()
    # remove common code fences
    text = re.sub(r"^```(?:json)?\s*", "", text)
    text = re.sub(r"\s*```$", "", text)

    # find first JSON object by matching braces to handle nested objects
    start = text.find("{")
    if start == -1:
        raise ValueError("No JSON object found in AI response.")

    depth = 0
    end = -1
    for i in range(start, len(text)):
        ch = text[i]
        if ch == '{':
            depth += 1
        elif ch == '}':
            depth -= 1
            if depth == 0:
                end = i
                break

    if end == -1:
        raise ValueError("No matching closing '}' found for JSON object.")

    json_text = text[start:end + 1]

    # common quick fixes: remove trailing commas before closing brackets/braces
    json_text = re.sub(r",\s*}", "}", json_text)
    json_text = re.sub(r",\s*]", "]", json_text)

    try:
        return json.loads(json_text)
    except json.JSONDecodeError as e:
        # make a best-effort additional cleanup and show a helpful error if parsing still fails
        fixed = json_text.replace("“", '"').replace("”", '"').replace("\n", " ")
        fixed = re.sub(r",\s*}", "}", fixed)
        fixed = re.sub(r",\s*]", "]", fixed)
        try:
            return json.loads(fixed)
        except json.JSONDecodeError:
            snippet = json_text[:1000] + ("..." if len(json_text) > 1000 else "")
            raise ValueError(f"Failed to parse JSON from AI response: {e.msg} (char {e.pos}).\nSnippet:\n{snippet}")


def add_textbox(slide, text, x, y, w, h, size=20, bold=False, color=(255, 255, 255)):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.text = text

    for paragraph in frame.paragraphs:
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = RGBColor(*color)

    return box


def add_bullets(slide, bullets, x, y, w, h):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()

    for bullet in bullets[:5]:
        p = frame.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(226, 232, 240)
        p.space_after = Pt(10)

    return box


def add_card(slide, x, y, w, h, fill_color=(30, 41, 59), line_color=(56, 189, 248)):
    shape = slide.shapes.add_shape(
        1,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill_color)
    shape.line.color.rgb = RGBColor(*line_color)
    return shape


def create_pptx(slide_content, output_path="research_presentation.pptx"):
    data = extract_json(slide_content)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slides_data = data.get("slides", [])[:8]

    for i, slide_data in enumerate(slides_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Background
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(15, 23, 42)

        title = slide_data.get("title", f"Slide {i + 1}")
        subtitle = slide_data.get("subtitle", "")
        bullets = slide_data.get("bullets", [])
        visual = slide_data.get("visual", "AI Visual Concept")

        # Title slide layout
        if i == 0:
            add_textbox(
                slide,
                title,
                0.8,
                1.3,
                11.8,
                1.2,
                size=44,
                bold=True,
                color=(255, 255, 255)
            )

            add_textbox(
                slide,
                subtitle,
                0.85,
                2.45,
                10.5,
                0.8,
                size=20,
                color=(148, 163, 184)
            )

            add_card(slide, 0.9, 4.1, 3.5, 1.4)
            add_textbox(slide, "RAG Powered", 1.2, 4.35, 3, 0.5, size=22, bold=True, color=(125, 211, 252))

            add_card(slide, 4.9, 4.1, 3.5, 1.4)
            add_textbox(slide, "AI Summary", 5.25, 4.35, 3, 0.5, size=22, bold=True, color=(125, 211, 252))

            add_card(slide, 8.9, 4.1, 3.5, 1.4)
            add_textbox(slide, "Research Insights", 9.1, 4.35, 3, 0.5, size=22, bold=True, color=(125, 211, 252))

        # Alternate split layout
        elif i % 2 == 1:
            add_textbox(slide, title, 0.7, 0.45, 12, 0.8, size=34, bold=True)
            add_textbox(slide, subtitle, 0.75, 1.2, 11.5, 0.5, size=15, color=(148, 163, 184))

            add_card(slide, 0.8, 2.0, 7.2, 4.6)
            add_bullets(slide, bullets, 1.1, 2.3, 6.6, 3.9)

            add_card(slide, 8.55, 2.0, 3.9, 4.6, fill_color=(8, 47, 73))
            add_textbox(slide, visual, 8.95, 3.25, 3.1, 1.3, size=20, bold=True, color=(186, 230, 253))

        # Card grid layout
        else:
            add_textbox(slide, title, 0.7, 0.45, 12, 0.8, size=34, bold=True)
            add_textbox(slide, subtitle, 0.75, 1.2, 11.5, 0.5, size=15, color=(148, 163, 184))

            positions = [
                (0.8, 2.0),
                (4.75, 2.0),
                (8.7, 2.0),
                (2.8, 4.55),
                (6.75, 4.55),
            ]

            for idx, bullet in enumerate(bullets[:5]):
                x, y = positions[idx]
                add_card(slide, x, y, 3.6, 1.55)
                add_textbox(slide, bullet, x + 0.25, y + 0.25, 3.1, 1.0, size=15, bold=False)

        # Footer
        add_textbox(
            slide,
            f"Research Paper Assistant | Slide {i + 1}",
            0.7,
            7.0,
            12,
            0.3,
            size=10,
            color=(148, 163, 184)
        )

    prs.save(output_path)
    return output_path